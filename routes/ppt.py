from flask import Blueprint, render_template, request, send_file, flash, redirect, url_for, current_app, session
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from utils.decorators import login_required
import os

ppt_bp = Blueprint("ppt", __name__)


@ppt_bp.route("/", methods=["GET", "POST"])
@login_required
def create_ppt():
    if request.method == "POST":
        title = request.form.get("title", "Presentation")
        slides_text = request.form.get("slides", "")

        if not slides_text.strip():
            flash("Please add slide content.", "danger")
            return redirect(url_for("ppt.create_ppt"))

        try:
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            # Title slide
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            if slide.placeholders[1]:
                slide.placeholders[1].text = "Department of CS & ML\nLoyola Academy"

            # Content slides
            slides_list = slides_text.strip().split("\n---\n")
            for slide_content in slides_list:
                lines = slide_content.strip().split("\n")
                if not lines:
                    continue

                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)

                slide_title = lines[0].strip().lstrip("#").strip()
                slide.shapes.title.text = slide_title

                body = "\n".join(lines[1:]).strip()
                if body and slide.placeholders[1]:
                    slide.placeholders[1].text = body

            os.makedirs(current_app.config["PPT_FOLDER"], exist_ok=True)
            filename = f"{title.replace(' ', '_')}.pptx"
            filepath = os.path.join(current_app.config["PPT_FOLDER"], filename)
            prs.save(filepath)

            return send_file(filepath, as_attachment=True, download_name=filename)

        except Exception as e:
            flash(f"Error creating PPT: {str(e)}", "danger")
            return redirect(url_for("ppt.create_ppt"))

    return render_template("ppt/create.html")